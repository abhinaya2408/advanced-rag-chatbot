from langchain_community.document_loaders import (
    PyPDFLoader,
    TextLoader,
    CSVLoader,
    BSHTMLLoader,
    Docx2txtLoader
)

from langchain.schema import Document

from pptx import Presentation


def load_pptx(file_path):

    prs = Presentation(file_path)

    text = ""

    for slide in prs.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                text += shape.text + "\n"

    return [
        Document(
            page_content=text,
            metadata={
                "source": file_path
            }
        )
    ]


def load_document(file_path):

    try:

        if file_path.endswith(".pdf"):

            loader = PyPDFLoader(file_path)

            return loader.load()

        elif file_path.endswith(".txt"):

            loader = TextLoader(file_path)

            return loader.load()

        elif file_path.endswith(".csv"):

            loader = CSVLoader(file_path)

            return loader.load()

        elif file_path.endswith(".html"):

            loader = BSHTMLLoader(file_path)

            return loader.load()

        elif file_path.endswith(".docx"):

            loader = Docx2txtLoader(file_path)

            return loader.load()

        elif file_path.endswith(".pptx"):

            return load_pptx(file_path)

        else:

            return []

    except Exception as e:

        print(
            f"Error loading {file_path}: {e}"
        )

        return []